"""
Build the Umubyeyi progress-presentation deck (.pptx).

Mirrors the academic capstone narrative arc the supervisor shared
(Intro -> Need -> Data -> Preprocessing/Exploration -> Approach ->
Results -> Discussion -> Conclusions), grounded in the project's real
datasets, EDA figures and trained-model baseline numbers.

Run:  python src/build_slides.py
Out:  docs/Umubyeyi_Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "reports", "figures")
OUT = os.path.join(ROOT, "docs", "Umubyeyi_Presentation.pptx")

# ---- palette (ties to the project's teal/coral figures) -------------------
TEAL = RGBColor(0x1A, 0x99, 0x88)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
LAV = RGBColor(0x9B, 0x8B, 0xC4)
DARK = RGBColor(0x2B, 0x2D, 0x42)
GRAY = RGBColor(0x55, 0x57, 0x6B)
LIGHT = RGBColor(0xF6, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def footer(slide, idx):
    rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), TEAL)
    tf = textbox(slide, Inches(0.3), SH - Inches(0.33), Inches(9), Inches(0.32),
                 MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run()
    r.text = "Umubyeyi  ·  Kinyarwanda Maternal Mental-Wellness Chatbot  ·  R. Irutingabo (ALU)"
    _set(r, 9, WHITE)
    tf2 = textbox(slide, SW - Inches(1.0), SH - Inches(0.33), Inches(0.7), Inches(0.32),
                  MSO_ANCHOR.MIDDLE)
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run()
    r2.text = str(idx)
    _set(r2, 9, WHITE)


def header(slide, kicker, title):
    rect(slide, 0, 0, Inches(0.22), SH, TEAL)  # left spine
    tf = textbox(slide, Inches(0.55), Inches(0.35), Inches(12), Inches(0.35))
    r = tf.paragraphs[0].add_run()
    r.text = kicker.upper()
    _set(r, 12, CORAL, bold=True)
    tf2 = textbox(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.9))
    r2 = tf2.paragraphs[0].add_run()
    r2.text = title
    _set(r2, 30, DARK, bold=True)
    rect(slide, Inches(0.55), Inches(1.62), Inches(1.4), Pt(3), CORAL)


def bullets(slide, items, x=Inches(0.6), y=Inches(1.95),
            w=Inches(7.4), h=Inches(4.9), size=17):
    tf = textbox(slide, x, y, w, h)
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(10)
        p.space_before = Pt(2)
        bullet = "•" if lvl == 0 else "–"
        r = p.add_run()
        r.text = f"{'    ' * lvl}{bullet}  {text}"
        _set(r, size if lvl == 0 else size - 2,
             DARK if lvl == 0 else GRAY, bold=(lvl == 0))


def image(slide, path, x, y, w, caption=None):
    if not os.path.exists(path):
        return
    pic = slide.shapes.add_picture(path, x, y, width=w)
    if caption:
        tf = textbox(slide, x, y + pic.height + Inches(0.05), w, Inches(0.45))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = caption
        _set(r, 10, GRAY, italic=True)


def new(kicker, title):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    return s


# ======================================================================
# 1. TITLE
# ======================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(2.55), SW, Inches(0.06), CORAL)
rect(s, 0, SH - Inches(0.5), SW, Inches(0.5), TEAL)
tf = textbox(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "BSc Software Engineering Capstone — Progress Review"
_set(r, 16, LAV, bold=True)
tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.1))
r = tf.paragraphs[0].add_run(); r.text = "Umubyeyi"
_set(r, 54, WHITE, bold=True)
tf = textbox(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.4))
r = tf.paragraphs[0].add_run()
r.text = "A Kinyarwanda-Language AI Chatbot for Maternal Mental Wellness in the Postpartum Window"
_set(r, 24, WHITE)
tf = textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2))
for line, col in [("Raissa Irutingabo  ·  African Leadership University", WHITE),
                  ("Supervisor: Samiratu Ntohsi", LAV),
                  ("Progress presentation — June 2026", GRAY)]:
    p = tf.paragraphs[0] if line.startswith("Raissa") else tf.add_paragraph()
    r = p.add_run(); r.text = line; _set(r, 15, col, bold=line.startswith("Raissa"))

# ======================================================================
# 2. OUTLINE  (mirrors the supervisor's example outline slide)
# ======================================================================
s = new("Agenda", "Outline")
header(s, "Agenda", "Outline")
bullets(s, [
    ("Introduction", 0),
    ("The Need", 0),
    ("Data", 0),
    ("Training corpus, framing dataset, response templates", 1),
    ("Data Preprocessing & Exploration", 0),
    ("The Approach", 0),
    ("System pipeline + the model ladder (floor → embeddings → transformer)", 1),
    ("Results", 0),
    ("Baseline & the cross-lingual degradation finding", 1),
    ("Discussion", 0),
    ("Conclusions & Next Steps", 0),
], y=Inches(1.9), w=Inches(11), size=18)
footer(s, 2)

# ======================================================================
# 3. INTRODUCTION
# ======================================================================
s = new("Introduction", "What is Umubyeyi?")
header(s, "Introduction", "What is Umubyeyi?")
bullets(s, [
    ("A browser-based chatbot that talks to first-time Rwandan mothers in "
     "Kinyarwanda about their mental wellbeing.", 0),
    ("Scope: the 0–6 month postpartum window — when low mood, anxiety and "
     "feeling overwhelmed peak.", 0),
    ("Informational and supportive — never diagnostic. Every reply carries a "
     "disclaimer and a safety pathway.", 0),
    ("Template-based, not generative: each intent maps to one guideline-sourced, "
     "human-reviewed response template (avoids LLM hallucination).", 0),
    ("Designed for low-resource deployment — classical ML, no GPU at "
     "inference, runs on Streamlit Community Cloud.", 0),
], w=Inches(7.3))
rect(s, Inches(8.2), Inches(2.1), Inches(4.5), Inches(4.4), LIGHT)
tf = textbox(s, Inches(8.45), Inches(2.35), Inches(4.0), Inches(4.0))
r = tf.paragraphs[0].add_run(); r.text = "In one sentence"
_set(r, 14, CORAL, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run()
r.text = ("“Meet mothers in their own language, route them to safe, "
          "expert-validated guidance, and measure exactly what translation "
          "costs us along the way.”")
_set(r, 16, DARK, italic=True)
footer(s, 3)

# ======================================================================
# 4. THE NEED
# ======================================================================
s = new("The Need", "Why this matters")
header(s, "The Need", "Why this matters")
bullets(s, [
    ("Perinatal mental-health conditions are common and under-served in "
     "sub-Saharan Africa, yet rarely screened in routine postpartum care.", 0),
    ("First-time mothers face the steepest adjustment, often in isolation and "
     "without a safe, private place to ask questions.", 0),
    ("Language is the barrier: almost all digital maternal-health tools are "
     "English-only.", 0),
    ("Kinyarwanda is severely under-resourced in NLP — no public mental-wellness "
     "intent dataset exists.", 0),
    ("Existing maternal-health chatbots are English-only and physical-health "
     "focused, and name translation as an unsolved limitation.", 0),
], w=Inches(11.5))
rect(s, Inches(0.6), Inches(5.75), Inches(11.5), Inches(0.7), CORAL)
tf = textbox(s, Inches(0.85), Inches(5.8), Inches(11), Inches(0.6), MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run()
r.text = ("The gap = Kinyarwanda  ×  mental wellness  ×  Rwanda  —  exactly where "
          "Umubyeyi sits.")
_set(r, 16, WHITE, bold=True)
footer(s, 4)

# ======================================================================
# 5. DATA
# ======================================================================
s = new("Data", "Two data sources, two distinct roles")
header(s, "Data", "Two data sources, two distinct roles")
cards = [
    ("Amod counseling Q&A", "TRAINING DATA", TEAL,
     ["3,512 rows of real therapist Q&A",
      "995 unique questions (heavy duplication)",
      "Filtered to maternal-relevant subset",
      "Labeled into 6 wellness intents"]),
    ("WHO / UNFPA guidance", "RESPONSE TEMPLATES", CORAL,
     ["Tier-1 global authorities",
      "Source of every user-facing answer",
      "Pre-translated + human-reviewed",
      "6 response templates (1 per intent)"]),
]
cw = Inches(3.95); gap = Inches(0.2); x0 = Inches(0.55); y0 = Inches(2.0)
for i, (title, tag, col, pts) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, Inches(4.4), LIGHT)
    rect(s, x, y0, cw, Inches(0.55), col)
    tf = textbox(s, x + Inches(0.15), y0 + Inches(0.03), cw - Inches(0.3), Inches(0.5),
                 MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run(); r.text = tag; _set(r, 11, WHITE, bold=True)
    tf = textbox(s, x + Inches(0.2), y0 + Inches(0.7), cw - Inches(0.4), Inches(0.7))
    r = tf.paragraphs[0].add_run(); r.text = title; _set(r, 16, DARK, bold=True)
    tf = textbox(s, x + Inches(0.2), y0 + Inches(1.45), cw - Inches(0.4), Inches(2.8))
    first = True
    for pt in pts:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(7)
        r = p.add_run(); r.text = "•  " + pt; _set(r, 13, GRAY)
footer(s, 5)

# ======================================================================
# 6. DATA PREPROCESSING & EXPLORATION
# ======================================================================
s = new("Data Preprocessing & Exploration", "What the data told us")
header(s, "Data Preprocessing & Exploration", "What the data told us")
bullets(s, [
    ("Deduplication finding: 3,512 rows collapse to 995 unique questions — "
     "2,517 are the same question answered by different counselors.", 0),
    ("Consequence: split on unique questions (GroupShuffleSplit) to avoid "
     "leakage and inflated accuracy.", 0),
    ("Intents are imbalanced — relationship/support dominates, self-care is "
     "scarce.", 0),
    ("This is why we report F1, not just accuracy.", 1),
    ("Cleaning: scope filtering, lowercasing, TF-IDF features.", 0),
], x=Inches(0.6), y=Inches(1.95), w=Inches(6.1), size=15)
image(s, os.path.join(FIG, "02_intent_coverage.png"),
      Inches(6.95), Inches(2.1), Inches(6.0),
      caption="Intent coverage across the maternal subset (class imbalance)")
footer(s, 6)

# ======================================================================
# 7. THE APPROACH
# ======================================================================
s = new("The Approach", "Part 1 — the system pipeline")
header(s, "The Approach (1/2)", "The system pipeline")
image(s, os.path.join(FIG, "nb_architecture.png"),
      Inches(0.7), Inches(2.0), Inches(12.0))
bullets(s, [
    ("Kinyarwanda in → NLLB-200 translates to English → TF-IDF + ML classifier "
     "assigns 1 of 6 intents → return that intent's response template → translate back → "
     "display the Kinyarwanda response with a disclaimer.", 0),
    ("Two safety gates: (1) confidence gate — low-confidence inputs fall back to "
     "general self-care; (2) danger-sign keywords bypass everything → Rwanda "
     "helpline escalation.", 0),
    ("Deployment MVP built: architecture diagram, web mockup, a Streamlit app "
     "and a FastAPI sketch (Swagger UI for testing).", 0),
], x=Inches(0.6), y=Inches(4.25), w=Inches(12.0), size=14)
footer(s, 7)

# ======================================================================
# 7b. THE APPROACH — MODEL LADDER  (maps to "Encoding Modalities")
# ======================================================================
s = new("The Approach", "Part 2 — the model ladder")
header(s, "The Approach (2/2)", "The model ladder: floor → embeddings → transformer")
tf = textbox(s, Inches(0.6), Inches(1.7), Inches(12.2), Inches(0.5))
r = tf.paragraphs[0].add_run()
r.text = ("Three classifiers on the SAME labels and SAME 80/20 split — each adds "
          "exactly one idea and must beat the one below it.")
_set(r, 15, GRAY, italic=True)
ladder = [
    ("1", "TF-IDF + Logistic Regression", "English",
     "FLOOR — simple, fast, interpretable; the currently deployed path.",
     "Done · 0.61 acc / 0.60 F1", TEAL),
    ("2", "FastText cc.rw.300 + MLP", "Kinyarwanda",
     "EMBEDDINGS — dense subword vectors capture meaning + Kinyarwanda "
     "morphology / out-of-vocabulary words.",
     "Coded · pending 4.5 GB vectors", LAV),
    ("3", "AfroXLMR fine-tune", "Kinyarwanda",
     "CEILING — contextual transformer pretrained on Kinyarwanda; classifies "
     "directly, no translate-then-classify hop.",
     "Coded · pending GPU run", CORAL),
]
ry = Inches(2.35); rh = Inches(1.25); rgap = Inches(0.18)
for n, title, lang, desc, status, col in ladder:
    rect(s, Inches(0.6), ry, Inches(12.2), rh, LIGHT)
    rect(s, Inches(0.6), ry, Inches(0.85), rh, col)
    tf = textbox(s, Inches(0.6), ry, Inches(0.85), rh, MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = n; _set(r, 30, WHITE, bold=True)
    tf = textbox(s, Inches(1.65), ry + Inches(0.1), Inches(7.4), rh - Inches(0.2))
    r = tf.paragraphs[0].add_run(); r.text = f"{title}   ({lang})"
    _set(r, 16, DARK, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(4)
    r = p.add_run(); r.text = desc; _set(r, 12, GRAY)
    tf = textbox(s, Inches(9.2), ry, Inches(3.45), rh, MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = status; _set(r, 13, col, bold=True)
    ry = ry + rh + rgap
footer(s, 8)

# ======================================================================
# 8. RESULTS — BASELINE
# ======================================================================
s = new("Results", "Baseline (Model 1)")
header(s, "Results", "Baseline — the floor model (Model 1)")
image(s, os.path.join(FIG, "baseline_confusion.png"),
      Inches(0.55), Inches(1.95), Inches(6.4),
      caption="Confusion matrix — TF-IDF + Logistic Regression, 6 intents")
rect(s, Inches(7.3), Inches(1.95), Inches(5.45), Inches(4.55), LIGHT)
tf = textbox(s, Inches(7.55), Inches(2.15), Inches(5.0), Inches(4.25))
r = tf.paragraphs[0].add_run(); r.text = "Floor model — established"
_set(r, 15, CORAL, bold=True)
for label, val in [("Accuracy", "0.61"), ("F1", "0.60"),
                   ("Precision", "0.63"), ("Recall", "0.59")]:
    p = tf.add_paragraph(); p.space_before = Pt(6)
    r = p.add_run(); r.text = f"{label}:  "; _set(r, 14, DARK, bold=True)
    r = p.add_run(); r.text = val; _set(r, 14, TEAL, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(12)
r = p.add_run(); r.text = "Reading the matrix"
_set(r, 13, CORAL, bold=True)
for line in ["Strong on relationship/support & sleep (clear cues).",
             "Confusion: overwhelmed ↔ relationship & sadness — overlapping "
             "language, the hard intents.",
             "F1 < accuracy → the minority intents are where the work is."]:
    p = tf.add_paragraph(); p.space_before = Pt(4)
    r = p.add_run(); r.text = "•  " + line; _set(r, 12, GRAY)
p = tf.add_paragraph(); p.space_before = Pt(12)
r = p.add_run()
r.text = ("Honest INITIAL baseline. Models 2 & 3 (FastText, AfroXLMR) "
          "are coded and pending compute.")
_set(r, 11, GRAY, italic=True)
footer(s, 9)

# ======================================================================
# 9. RESULTS — DEGRADATION / SCIENTIFIC CONTRIBUTION
# ======================================================================
s = new("Results", "The cross-lingual degradation finding")
header(s, "Results", "The core scientific contribution")
bullets(s, [
    ("Central research question: how much accuracy do we lose to translation?", 0),
    ("Method: classify the English test set directly vs. after an "
     "EN→Kinyarwanda→EN round-trip through NLLB-200 — the gap is the "
     "degradation coefficient.", 0),
    ("This would be the first published Kinyarwanda mental-wellness "
     "classification benchmark.", 0),
], x=Inches(0.6), y=Inches(1.95), w=Inches(7.0), size=15)
rect(s, Inches(7.9), Inches(2.0), Inches(4.9), Inches(4.4), DARK)
tf = textbox(s, Inches(8.15), Inches(2.2), Inches(4.4), Inches(4.1))
r = tf.paragraphs[0].add_run(); r.text = "Live demo finding"
_set(r, 14, CORAL, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run(); r.text = "Input (Kinyarwanda):"
_set(r, 12, LAV, bold=True)
p = tf.add_paragraph()
r = p.add_run(); r.text = "“umugabo wanjye ntamfasha”  (my husband does NOT help me)"
_set(r, 13, WHITE, italic=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run(); r.text = "NLLB dropped the negation →"
_set(r, 12, LAV, bold=True)
p = tf.add_paragraph()
r = p.add_run(); r.text = "“my husband is helping me”  — meaning inverted"
_set(r, 13, WHITE, italic=True)
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run()
r.text = ("Confidence fell to 0.31 → the safety gate caught it → safe fallback. "
          "Live proof of both the risk AND the mitigation.")
_set(r, 13, WHITE)
footer(s, 10)

# ======================================================================
# 10. DISCUSSION
# ======================================================================
s = new("Discussion", "What we are learning")
header(s, "Discussion", "Discussion")
bullets(s, [
    ("Translation hallucination is the real risk — not classifier accuracy. "
     "Dropped negations can invert meaning.", 0),
    ("The two-gate design (calibrated confidence + danger-sign override) turns that risk "
     "into a controlled, safe fallback rather than a wrong answer.", 0),
    ("Open finding: the proposed fixed 0.40 confidence threshold abstains ~88% of "
     "the time — it needs empirical tuning, not assumption.", 0),
    ("“Strong ML” here = rigor + a genuinely novel measurement, not model "
     "size. Classical models are the right call for low-resource deployment.", 0),
], w=Inches(11.8), size=16)
footer(s, 11)

# ======================================================================
# 11. CONCLUSIONS & NEXT STEPS
# ======================================================================
s = new("Conclusions", "Conclusions & next steps")
header(s, "Conclusions", "Conclusions & Next Steps")
tf = textbox(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "Where we are"
_set(r, 16, TEAL, bold=True)
bullets(s, [
    ("End-to-end Kinyarwanda pipeline works (NLLB + classifier + response templates + safety).", 0),
    ("Floor baseline established: 0.61 acc / 0.60 F1 (Model 1).", 0),
    ("Principled model ladder designed; embedding + transformer arms coded.", 0),
    ("Translation-risk demonstrated live, and contained by the safety gate.", 0),
], x=Inches(0.6), y=Inches(2.3), w=Inches(6.0), size=14)
tf = textbox(s, Inches(6.9), Inches(1.85), Inches(6.0), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "Next steps"
_set(r, 16, CORAL, bold=True)
bullets(s, [
    ("Measure the cross-lingual degradation coefficient.", 0),
    ("AfroXLMR fine-tune as a comparison arm.", 0),
    ("Future: build a true evidence-sourced retrieval KB (beyond the 6 templates).", 0),
    ("SUS usability study with 10 proxy users (target > 70).", 0),
], x=Inches(6.9), y=Inches(2.3), w=Inches(6.0), size=14)
rect(s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.65), TEAL)
tf = textbox(s, Inches(0.85), Inches(5.9), Inches(11.7), Inches(0.55), MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run()
r.text = ("Contribution: the first Kinyarwanda maternal mental-wellness classifier "
          "+ a quantified benchmark of what translation costs.")
_set(r, 15, WHITE, bold=True)
footer(s, 12)

# ======================================================================
# 12. THANK YOU
# ======================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(3.4), SW, Inches(0.06), CORAL)
tf = textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2))
r = tf.paragraphs[0].add_run(); r.text = "Murakoze — Thank you"
_set(r, 44, WHITE, bold=True)
tf = textbox(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.8))
r = tf.paragraphs[0].add_run(); r.text = "Questions & discussion"
_set(r, 22, LAV)
tf = textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6))
r = tf.paragraphs[0].add_run()
r.text = "Raissa Irutingabo  ·  github.com/IrutingaboRaissa/UMUBYEYI"
_set(r, 14, GRAY)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
try:
    prs.save(OUT)
    dest = OUT
except PermissionError:
    dest = OUT.replace(".pptx", "_NEW.pptx")
    prs.save(dest)
    print("NOTE: original file was open/locked — saved to a new file instead.")
print("Saved:", dest, "·", len(prs.slides._sldIdLst), "slides")
